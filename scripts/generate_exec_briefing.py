# SPDX-License-Identifier: MIT
# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx>=1.0.2"]
# ///
"""Generate the one-slide Halcyon executive briefing."""

from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.slide import Slide
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "presentations" / "2026-08-15-exec-briefing.pptx"

BACKGROUND = RGBColor(17, 20, 24)
SURFACE = RGBColor(33, 37, 43)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(190, 197, 206)
AMD_GOLD = RGBColor(193, 169, 104)
AMD_TEAL = RGBColor(0, 194, 222)
AMD_RED = RGBColor(237, 28, 36)


def add_text(
    slide: Slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int,
    color: RGBColor = WHITE,
    bold: bool = False,
) -> None:
    """Add a consistently styled text box to a slide."""
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_panel(
    slide: Slide,
    *,
    left: float,
    heading: str,
    heading_color: RGBColor,
    bullets: list[str],
) -> None:
    """Add one executive-briefing content panel."""
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(2.25),
        Inches(4.05),
        Inches(4.35),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = SURFACE
    panel.line.color.rgb = RGBColor(58, 64, 72)

    add_text(
        slide,
        left=left + 0.18,
        top=2.42,
        width=3.68,
        height=0.45,
        text=heading,
        size=18,
        color=heading_color,
        bold=True,
    )
    body = "\n\n".join(f"• {item}" for item in bullets)
    add_text(
        slide,
        left=left + 0.18,
        top=2.95,
        width=3.68,
        height=3.35,
        text=body,
        size=14,
        color=WHITE,
    )


def build() -> None:
    """Build and validate the executive-briefing PowerPoint."""
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BACKGROUND

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        presentation.slide_width,
        Inches(0.1),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = AMD_RED
    accent.line.fill.background()

    add_text(
        slide,
        left=0.45,
        top=0.28,
        width=12.4,
        height=0.58,
        text="Halcyon AI Infrastructure — Executive Briefing | 2026-08-15",
        size=25,
        bold=True,
    )
    add_text(
        slide,
        left=0.45,
        top=0.98,
        width=12.4,
        height=0.9,
        text=(
            "AT RISK — 5/5 primary requirements mapped; "
            "7/7 production-readiness gates remain FAIL"
        ),
        size=20,
        color=AMD_RED,
        bold=True,
    )

    add_panel(
        slide,
        left=0.45,
        heading="Key Wins",
        heading_color=AMD_GOLD,
        bullets=[
            "Selected one managed, low-operations Part 1 stack.",
            "Added vendor isolation, upload security, and recoverable queue semantics.",
        ],
    )
    add_panel(
        slide,
        left=4.64,
        heading="Risks & Decisions",
        heading_color=AMD_RED,
        bullets=[
            "RPO/RTO, budget, traffic, and incident forensics remain unresolved.",
            "Identity, quarantine/scanning, retention, and encryption need owners by Aug 22.",
        ],
    )
    add_panel(
        slide,
        left=8.83,
        heading="Next 2 Weeks",
        heading_color=AMD_TEAL,
        bullets=[
            "Approve ADR-001 and close client decisions.",
            "Implement secure async simulation plus Terraform.",
            "Prove isolation, retries, rollback, restore, load, and headroom.",
        ],
    )

    add_text(
        slide,
        left=0.48,
        top=6.82,
        width=12.2,
        height=0.35,
        text=(
            "Executive ask: approve Option B as proposed; do not claim production "
            "readiness until every evidence gate passes."
        ),
        size=12,
        color=MUTED,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(OUTPUT)
    if not zipfile.is_zipfile(OUTPUT):
        raise RuntimeError(f"Generated file is not a valid PPTX archive: {OUTPUT}")
    check = Presentation(OUTPUT)
    if len(check.slides) != 1:
        raise RuntimeError("Executive briefing must contain exactly one slide.")
    print(OUTPUT)


if __name__ == "__main__":
    build()
